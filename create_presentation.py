"""
PowerPoint Presentation Generator for Hospital Readmission Prediction Project
Creates a comprehensive presentation for project submission
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style the title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Add content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(18)
        p.level = 0
        p.space_after = Pt(12)
        
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, point in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.level = 0
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, point in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.level = 0
        p.space_after = Pt(8)
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add slide with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add image if it exists
    if os.path.exists(image_path):
        left = Inches(1.5)
        top = Inches(1.2)
        height = Inches(5)
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        # Add placeholder text
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = f"[Image: {os.path.basename(image_path)}]"
        p = placeholder_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    
    # Add caption
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        p = caption_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Predicting Hospital Readmission Risk",
        "A Comparative Study of Classical and Ensemble Models\n\n" +
        "Adele Chinda, Oumar Diallo, Yusuf Mumin\n" +
        "Georgia State University\nData Mining Project\nNovember 2024"
    )
    
    # Slide 2: Problem Statement
    add_content_slide(
        prs,
        "Problem Statement",
        [
            "Hospital readmissions are costly and often preventable",
            "30-day readmission rates are a key quality metric for healthcare",
            "Challenge: Predict which patients are at high risk of readmission",
            "Goal: Enable early intervention and reduce readmission rates",
            "Dataset: 100,000+ hospital admissions from diabetes patients (1999-2008)"
        ]
    )
    
    # Slide 3: Research Objectives
    add_content_slide(
        prs,
        "Research Objectives",
        [
            "Build a reproducible ML pipeline for readmission prediction",
            "Compare multiple modeling approaches:",
            "   • Logistic Regression (interpretable baseline)",
            "   • Multilayer Perceptron (nonlinear neural network)",
            "   • XGBoost (ensemble tree boosting)",
            "Address class imbalance using SMOTE",
            "Evaluate trade-offs between interpretability and performance",
            "Identify key predictive features for clinical insight"
        ]
    )
    
    # Slide 4: Dataset Overview
    add_two_column_slide(
        prs,
        "Dataset Overview",
        [
            "📊 Source:",
            "UCI Diabetes 130-US Hospitals",
            "(1999-2008)",
            "",
            "📈 Size:",
            "~100,000 hospital admissions",
            "50+ features",
            "",
            "🎯 Target Variable:",
            "Readmitted within 30 days",
            "(Binary classification)",
            "",
            "⚖️ Class Imbalance:",
            "~11% positive cases"
        ],
        [
            "📋 Feature Categories:",
            "",
            "• Demographics",
            "  - Age, gender, race",
            "",
            "• Clinical History",
            "  - Diagnoses (ICD-9)",
            "  - Comorbidities",
            "",
            "• Visit Details",
            "  - Time in hospital",
            "  - Number of procedures",
            "  - Lab tests performed",
            "",
            "• Medications",
            "  - Drug prescriptions",
            "  - Insulin usage"
        ]
    )
    
    # Slide 5: Data Pipeline
    add_content_slide(
        prs,
        "Data Processing Pipeline",
        [
            "1. Data Cleaning",
            "   • Handle missing values (imputation and removal)",
            "   • Remove low-variance and irrelevant features",
            "",
            "2. Feature Engineering",
            "   • Encode categorical variables (one-hot, label encoding)",
            "   • Create derived features (total visits, medication counts)",
            "   • Scale numerical features (StandardScaler)",
            "",
            "3. Feature Selection",
            "   • XGBoost feature importance analysis",
            "   • Select top 100 most predictive features",
            "",
            "4. Imbalance Handling",
            "   • Apply SMOTE (Synthetic Minority Over-sampling Technique)"
        ]
    )
    
    # Slide 6: Methodology - Models
    add_two_column_slide(
        prs,
        "Modeling Approaches",
        [
            "🔵 Logistic Regression",
            "• Linear model with L2 regularization",
            "• Interpretable coefficients",
            "• Baseline performance",
            "• Log-odds ratios for features",
            "",
            "🟢 Multilayer Perceptron",
            "• 2 hidden layers (64, 32 units)",
            "• ReLU activation",
            "• Captures nonlinear patterns",
            "• Early stopping for regularization"
        ],
        [
            "🟠 XGBoost",
            "• Gradient boosted decision trees",
            "• Handles missing values natively",
            "• Built-in feature importance",
            "• Class weighting for imbalance",
            "• Regularization parameters",
            "",
            "⚙️ Training Strategy",
            "• 80/20 train-test split",
            "• Cross-validation",
            "• Hyperparameter tuning",
            "• Threshold optimization"
        ]
    )
    
    # Slide 7: Mathematical Formulation
    add_content_slide(
        prs,
        "Model Formulations",
        [
            "Logistic Regression:",
            "   P(Y=1|x) = σ(w^T x + b)",
            "   Loss = -1/N Σ[y log ŷ + (1-y)log(1-ŷ)] + λ||w||²",
            "",
            "MLP Architecture:",
            "   h₁ = ReLU(W₁x + b₁)",
            "   h₂ = ReLU(W₂h₁ + b₂)",
            "   ŷ = σ(W₃h₂ + b₃)",
            "",
            "XGBoost Objective:",
            "   L = Σ l(yᵢ, ŷᵢ) + Σ Ω(fₖ)",
            "   Ω(f) = γT + ½λ||w||²"
        ]
    )
    
    # Slide 8: SMOTE Technique
    add_content_slide(
        prs,
        "Handling Class Imbalance: SMOTE",
        [
            "Challenge: Only 11% of admissions result in readmission",
            "",
            "SMOTE (Synthetic Minority Over-sampling Technique):",
            "   • Creates synthetic minority class samples",
            "   • Interpolates between existing minority instances",
            "   • x_new = xᵢ + α(xⱼ - xᵢ), where α ~ U(0,1)",
            "",
            "Benefits:",
            "   • Balances training data",
            "   • Improves model sensitivity to minority class",
            "   • Prevents overfitting compared to simple oversampling",
            "",
            "Implementation:",
            "   • Applied only to training set",
            "   • Test set remains unchanged for fair evaluation"
        ]
    )
    
    # Slide 9: Evaluation Metrics
    add_two_column_slide(
        prs,
        "Evaluation Metrics",
        [
            "📊 Primary Metrics:",
            "",
            "• ROC AUC",
            "  Discriminative ability",
            "",
            "• Precision-Recall AUC",
            "  Performance on imbalanced data",
            "",
            "• Precision",
            "  Accuracy of positive predictions",
            "",
            "• Recall (Sensitivity)",
            "  Coverage of actual positives",
            "",
            "• F1 Score",
            "  Harmonic mean of precision/recall"
        ],
        [
            "🎯 Why Multiple Metrics?",
            "",
            "• ROC AUC can be misleading",
            "  on imbalanced datasets",
            "",
            "• PR AUC better reflects",
            "  minority class performance",
            "",
            "• Clinical context matters:",
            "  - High recall: Don't miss at-risk",
            "    patients (false negatives costly)",
            "  - High precision: Avoid",
            "    unnecessary interventions",
            "    (false positives wasteful)",
            "",
            "• Threshold tuning allows",
            "  customization to priorities"
        ]
    )
    
    # Slide 10: Results - Performance Table
    add_content_slide(
        prs,
        "Model Performance Comparison",
        [
            "┌─────────────────────┬─────────┬─────────┬───────────┬────────┬──────┐",
            "│ Model               │ ROC AUC │ PR AUC  │ Precision │ Recall │  F1  │",
            "├─────────────────────┼─────────┼─────────┼───────────┼────────┼──────┤",
            "│ Logistic Regression │  0.56   │  0.35   │   0.89    │  0.97  │ 0.93 │",
            "├─────────────────────┼─────────┼─────────┼───────────┼────────┼──────┤",
            "│ MLP                 │  0.59   │  0.38   │   0.89    │  0.98  │ 0.93 │",
            "├─────────────────────┼─────────┼─────────┼───────────┼────────┼──────┤",
            "│ XGBoost             │  0.62   │  0.42   │   0.93    │  0.40  │ 0.56 │",
            "└─────────────────────┴─────────┴─────────┴───────────┴────────┴──────┘",
            "",
            "Key Findings:",
            "✓ XGBoost: Best discriminative performance (highest AUC)",
            "✓ LogReg/MLP: Best balance (highest F1 score)",
            "✓ Clear precision-recall trade-off across models"
        ]
    )
    
    # Slide 11: Results - ROC Curve
    add_image_slide(
        prs,
        "ROC Curve Comparison",
        "output/train/models_comparison_roc.png",
        "XGBoost achieves the highest ROC AUC (0.62), demonstrating superior discriminative ability"
    )
    
    # Slide 12: Feature Importance
    add_image_slide(
        prs,
        "Top Predictive Features",
        "output/logreg_top30_coefficients.png",
        "Logistic regression coefficients reveal key risk factors for readmission"
    )
    
    # Slide 13: Key Insights
    add_content_slide(
        prs,
        "Key Insights & Findings",
        [
            "1. Model Trade-offs:",
            "   • XGBoost excels at identifying high-risk patients (high precision)",
            "   • MLP/LogReg better at not missing at-risk patients (high recall)",
            "",
            "2. Feature Insights:",
            "   • Number of prior visits strongly predictive",
            "   • Medication changes indicate severity",
            "   • Lab procedures correlate with complications",
            "",
            "3. Class Imbalance Impact:",
            "   • SMOTE significantly improved minority class detection",
            "   • PR AUC more informative than ROC AUC",
            "",
            "4. Clinical Implications:",
            "   • Model choice depends on intervention cost vs. missed case cost",
            "   • Interpretable features enable clinician trust"
        ]
    )
    
    # Slide 14: Model Interpretation
    add_two_column_slide(
        prs,
        "Model Interpretation",
        [
            "🔍 Logistic Regression",
            "Strengths:",
            "• Direct coefficient interpretation",
            "• Clinician-friendly explanations",
            "• Fast inference",
            "",
            "Limitations:",
            "• Assumes linear relationships",
            "• May miss complex interactions",
            "",
            "🧠 Neural Network (MLP)",
            "Strengths:",
            "• Captures nonlinear patterns",
            "• Flexible architecture",
            "",
            "Limitations:",
            "• \"Black box\" nature",
            "• Requires more data",
            "• Longer training time"
        ],
        [
            "🌳 XGBoost",
            "Strengths:",
            "• Excellent predictive performance",
            "• Feature importance built-in",
            "• Handles missing data",
            "• Robust to outliers",
            "",
            "Limitations:",
            "• Less interpretable than LogReg",
            "• More hyperparameters to tune",
            "• Computationally intensive",
            "",
            "💡 Recommendation:",
            "• Use XGBoost for prediction",
            "• Use LogReg for interpretation",
            "• Ensemble approaches possible"
        ]
    )
    
    # Slide 15: Challenges Faced
    add_content_slide(
        prs,
        "Challenges & Solutions",
        [
            "Challenge 1: Severe Class Imbalance (11% positive)",
            "   → Solution: SMOTE + class weighting + threshold tuning",
            "",
            "Challenge 2: High-Dimensional Feature Space (100+ features)",
            "   → Solution: Feature selection using XGBoost importance scores",
            "",
            "Challenge 3: Missing Data",
            "   → Solution: Strategic imputation + removal of sparse features",
            "",
            "Challenge 4: Model Interpretability vs. Performance",
            "   → Solution: Compare multiple model paradigms, document trade-offs",
            "",
            "Challenge 5: Evaluation Metric Selection",
            "   → Solution: Report multiple metrics (ROC AUC, PR AUC, F1)"
        ]
    )
    
    # Slide 16: Future Directions
    add_content_slide(
        prs,
        "Future Work & Extensions",
        [
            "1. Temporal Modeling",
            "   • RNNs/LSTMs for sequential patient histories",
            "   • Capture temporal patterns in readmission risk",
            "",
            "2. Advanced Techniques",
            "   • Deep learning: Transformers for EHR data",
            "   • Ensemble methods: Stack multiple models",
            "   • Calibration: Platt scaling, isotonic regression",
            "",
            "3. Causal Analysis",
            "   • Move beyond correlation to causation",
            "   • Propensity score matching, instrumental variables",
            "",
            "4. Deployment Considerations",
            "   • Fairness across demographic groups",
            "   • Real-time prediction API",
            "   • Integration with hospital EHR systems",
            "   • Continuous model monitoring and retraining"
        ]
    )
    
    # Slide 17: Technical Implementation
    add_two_column_slide(
        prs,
        "Technical Stack & Tools",
        [
            "🐍 Programming:",
            "• Python 3.8+",
            "• Jupyter Notebooks",
            "",
            "📚 ML Libraries:",
            "• scikit-learn",
            "• XGBoost",
            "• imbalanced-learn (SMOTE)",
            "",
            "📊 Data Processing:",
            "• pandas",
            "• NumPy",
            "",
            "📈 Visualization:",
            "• matplotlib",
            "• seaborn",
            "• plotly"
        ],
        [
            "🗂️ Project Structure:",
            "• Modular pipeline design",
            "• Separate scripts for each stage",
            "• Reproducible workflows",
            "",
            "💾 Model Persistence:",
            "• Pickle files for trained models",
            "• Saved preprocessors & scalers",
            "",
            "📝 Documentation:",
            "• README with setup instructions",
            "• Code comments & docstrings",
            "• Project website (GitHub Pages)",
            "",
            "🔧 Version Control:",
            "• Git + GitHub",
            "• Clear commit history"
        ]
    )
    
    # Slide 18: Reproducibility
    add_content_slide(
        prs,
        "Reproducibility & Open Science",
        [
            "✅ Complete Pipeline Available:",
            "   • All code on GitHub: github.com/arrdel/patient-readmission-prediction",
            "   • requirements.txt for dependencies",
            "   • Step-by-step execution instructions",
            "",
            "✅ Dataset Accessibility:",
            "   • Public UCI ML Repository dataset",
            "   • Clear data provenance documentation",
            "",
            "✅ Experiment Tracking:",
            "   • Saved model checkpoints",
            "   • Performance metrics logged",
            "   • Visualizations for all experiments",
            "",
            "✅ Documentation:",
            "   • Comprehensive README",
            "   • Project website with interactive results",
            "   • Technical report with methodology details"
        ]
    )
    
    # Slide 19: Lessons Learned
    add_content_slide(
        prs,
        "Lessons Learned",
        [
            "1. Domain Knowledge Matters",
            "   • Understanding healthcare context crucial for feature engineering",
            "   • Clinical input would improve model utility",
            "",
            "2. Evaluation is Non-Trivial",
            "   • Single metric insufficient for imbalanced problems",
            "   • Threshold selection depends on deployment context",
            "",
            "3. Interpretability is Valuable",
            "   • Black-box models may perform better but harder to trust",
            "   • Simpler models often \"good enough\" and more deployable",
            "",
            "4. Data Quality > Fancy Algorithms",
            "   • Feature engineering had major impact",
            "   • Garbage in, garbage out principle holds",
            "",
            "5. Reproducibility Requires Effort",
            "   • Documentation and organization take time but pay dividends"
        ]
    )
    
    # Slide 20: Conclusions
    add_content_slide(
        prs,
        "Conclusions",
        [
            "✓ Successfully built end-to-end ML pipeline for readmission prediction",
            "",
            "✓ Compared three distinct modeling paradigms:",
            "   • Logistic Regression: Most interpretable",
            "   • MLP: Good balance of performance and complexity",
            "   • XGBoost: Best discriminative performance",
            "",
            "✓ Demonstrated effective handling of class imbalance via SMOTE",
            "",
            "✓ Identified key predictive features with clinical relevance",
            "",
            "✓ Highlighted practical trade-offs in model selection",
            "",
            "✓ Created fully reproducible, well-documented codebase",
            "",
            "➡️ Ready for further research and potential deployment"
        ]
    )
    
    # Slide 21: References
    add_content_slide(
        prs,
        "References & Resources",
        [
            "Dataset:",
            "• Strack et al. (2014). UCI Diabetes 130-US Hospitals Dataset",
            "  archive.ics.uci.edu/dataset/296",
            "",
            "Methods:",
            "• Chawla et al. (2002). SMOTE: Synthetic Minority Over-sampling Technique",
            "• Chen & Guestrin (2016). XGBoost: A Scalable Tree Boosting System",
            "• Pedregosa et al. (2011). Scikit-learn: Machine Learning in Python",
            "",
            "Project Resources:",
            "• Code: github.com/arrdel/patient-readmission-prediction",
            "• Website: [GitHub Pages URL]",
            "• Report: See repository docs/",
            "",
            "Contact:",
            "• Adele Chinda, Oumar Diallo, Yusuf Mumin",
            "• Georgia State University"
        ]
    )
    
    # Slide 22: Thank You
    add_title_slide(
        prs,
        "Thank You!",
        "Questions?\n\n" +
        "Project Repository:\ngithub.com/arrdel/patient-readmission-prediction\n\n" +
        "Adele Chinda • Oumar Diallo • Yusuf Mumin\n" +
        "Georgia State University"
    )
    
    return prs

def main():
    """Main function to create and save presentation"""
    print("Creating PowerPoint presentation...")
    
    prs = create_presentation()
    
    # Save presentation
    output_path = "Hospital_Readmission_Prediction_Presentation.pptx"
    prs.save(output_path)
    
    print(f"✓ Presentation created successfully: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
    print("\nSlide Overview:")
    slide_titles = [
        "Title Slide",
        "Problem Statement",
        "Research Objectives",
        "Dataset Overview",
        "Data Processing Pipeline",
        "Modeling Approaches",
        "Mathematical Formulation",
        "SMOTE Technique",
        "Evaluation Metrics",
        "Model Performance Comparison",
        "ROC Curve Comparison",
        "Top Predictive Features",
        "Key Insights & Findings",
        "Model Interpretation",
        "Challenges & Solutions",
        "Future Work & Extensions",
        "Technical Stack & Tools",
        "Reproducibility & Open Science",
        "Lessons Learned",
        "Conclusions",
        "References & Resources",
        "Thank You"
    ]
    
    for i, title in enumerate(slide_titles, 1):
        print(f"  {i}. {title}")
    
    print("\n📊 Presentation is ready for your project submission!")
    print("🖼️  Note: Image slides will show placeholder text if images are not found.")
    print("    Copy images from output/ folder to make them appear.")

if __name__ == "__main__":
    main()
